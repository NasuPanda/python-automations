import os
from copy import deepcopy

from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide, SlideLayout, SlideShapes

from ppt import types


class PresentationAccessor:
    """Power Point Accessor
    NOTE: slideのインデックス指定は0スタート.
    """

    def __init__(
        self,
        src_filepath: str,
        first_active_slide: types.SlideKey = 0,
        slide_layout: types.SlideKey = 0,
    ) -> None:
        """初期化処理"""
        self._filepath = src_filepath
        self._presentation: PresentationType
        self._active_slide: Slide
        self._slide_layout: SlideLayout
        self._load_presentation(first_active_slide, slide_layout)

    @staticmethod
    def emu_to_cm(emu_val: int) -> float:
        """
        emu→cmに変換 (1cm = 360000emu)
        """
        return emu_val / 360000

    @staticmethod
    def cm_to_emu(cm_val: int) -> int:
        """
        cm→emuに変換 (1cm = 360000emu)
        """
        return cm_val * 360000

    # Private
    def _load_presentation(self, active_slide: types.SlideKey, slide_layout: types.SlideKey) -> None:
        """private: pptxファイルのロード"""
        if os.path.isfile(self._filepath):
            self._presentation = Presentation(self._filepath)
        else:
            self._presentation = Presentation()
        self._active_slide = self._get_slide(active_slide)
        self._slide_layout = self._get_slide_layout(slide_layout)

    def _get_slide(self, slide_key: types.SlideKey) -> Slide:
        """private: スライド取得"""
        return self._presentation.slides[slide_key]

    def _get_slide_layout(self, slide_key: types.SlideKey) -> SlideLayout:
        """private: スライドレイアウト取得"""
        return self._presentation.slide_layouts[slide_key]

    def _get_slide_index(self, slide) -> int:
        """private: スライドのインデックスを取得する"""
        return self._presentation.slides.index(slide)

    @property
    def _shapes(self) -> SlideShapes:
        """private: アクティブなスライドのshapesを取得"""
        return self._active_slide.shapes

    # Public
    @property
    def slide_layouts(self) -> list[types.LayoutInfo]:
        """全てのスライドレイアウトのインデックス、名前を取得する"""
        layouts = self._presentation.slide_layouts
        return [types.LayoutInfo(name=layout.name, number=i) for i, layout in enumerate(layouts)]

    @property
    def number_of_slide(self) -> int:
        """スライドの合計数"""
        return len(self._presentation.slides)

    def get_contents(self, shape_type: int) -> list[types.Content]:
        """特定のshape_typeのshapeの情報を取得する
        See Also:
            types.Content
        """
        contents: list[types.Content] = []
        for shape in self._shapes:
            if shape.shape_type == shape_type:
                contents.append(types.Content(
                    coordinates=types.Coordinates(left=shape.left, top=shape.top),
                    size=types.Size(width=shape.width, height=shape.height),
                ))

        return [
            types.Content(
                coordinates=types.Coordinates(left=shape.left, top=shape.top),
                size=types.Size(width=shape.width, height=shape.height),
                text=shape.text if shape.text else ""
            )
            for shape in self._shapes
            if shape.shape_type == shape_type
        ]

    def change_active_slide(self, slide: types.SlideKey | Slide) -> None:
        """アクティブなスライドを変更する"""
        if isinstance(slide, int):
            self._active_slide = self._get_slide(slide)
        else:
            self._active_slide = slide

    def change_slide_layout(self, slide_key: types.SlideKey) -> None:
        """使用するスライドレイアウトを変更する"""
        slide_layout = self._presentation.slide_layouts[slide_key]
        self._slide_layout = slide_layout
        print(f"変更. 名前:{slide_layout.name}, インデックス:{slide_key}")

    def overwrite(self) -> bool:
        """上書き保存"""
        try:
            self._presentation.save(self._filepath)
        # 保存成功/失敗
        except PermissionError:
            return False
        return True

    def save_as(self, dst_filepath: str) -> bool:
        """名前を付けて保存"""
        try:
            self._presentation.save(dst_filepath)
        # 保存成功/失敗
        except PermissionError:
            return False
        return True

    def add_textbox(self, textbox_info: types.Content):
        """アクティブなスライドにテキストボックスを追加する"""
        textbox = self._shapes.add_textbox(textbox_info.coordinates.left, textbox_info.coordinates.top, textbox_info.size.width, textbox_info.size.height)
        textbox.text_frame.text = textbox_info.text

    def add_picture(self, picture_info: types.Picture):
        """アクティブなスライドに画像を追加する"""
        self._shapes.add_picture(picture_info.path, picture_info.coordinates.left, picture_info.coordinates.top, picture_info.size.width, picture_info.size.height)

    def add_slide(self, change_active_slide: bool = True) -> Slide:
        """self._slide_layout を使用してスライドを追加する"""
        slide = self._presentation.slides.add_slide(self._slide_layout)
        if change_active_slide:
            self.change_active_slide(slide)
        return slide

    def duplicate_slide(self, src_slide_index=0, change_active_slide: bool = True) -> Slide:
        """スライドを複製する
        参考 : https://stackoverflow.com/questions/50866634/python-pptx-copy-slide
        """
        src_slide = self._get_slide(src_slide_index)
        dist_slide = self.add_slide()

        for shape in src_slide.shapes:
            el = shape.element
            new_el = deepcopy(el)
            dist_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

        if change_active_slide:
            self.change_active_slide(dist_slide)
        return dist_slide

    def remove_slide(self, slide: types.SlideKey | Slide) -> None:
        """スライドを削除する
        NOTE: アクティブなスライドは削除出来ない(削除後の処理が面倒なので)
        """
        slide_index = slide if isinstance(slide, int) else self._get_slide_index(slide)
        if slide_index == self._get_slide_index(self._active_slide):
            print("[Error] Can't remove the active slide.")
            return

        xml_slides = self._presentation.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[slide_index])

    def clear_shape(self, shape_type: int) -> None:
        """アクティブなスライドから特定の shape_type を **全て** 削除する
        参考 : https://stackoverflow.com/questions/64700638/is-there-a-way-to-delete-a-shape-with-python-pptx
        """
        for shape in self._shapes:
            if shape.shape_type == shape_type:
                xml_reference = shape._sp
                xml_reference.getparent().remove(xml_reference)
