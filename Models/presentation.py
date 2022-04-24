from copy import deepcopy

from pptx import Presentation

import config
from Models.slide import Slide
from Models.content import Image, TextBox


class PresentationReader():
    """プレゼンテーションを読み込むクラス。
    """
    def __init__(self, path: str) -> None:
        self.prs: Presentation = Presentation(path)

    def get_imgs(self, slide_index=0) -> list[Image]:
        """スライドが持つ画像の情報を取得(RECTANGLEを使用)"""
        return self.__get_contents(
            content_type=config.IMAGE_KEY,
            shape_type=config.SHAPES["image"],
            slide_index=slide_index
        )

    def get_textbox(self, slide_index=0) -> list[TextBox]:
        """スライドが持つテキストボックスの情報を取得"""
        return self.__get_contents(
            content_type=config.TEXTBOX_KEY,
            shape_type=config.SHAPES["textbox"],
            slide_index=slide_index
        )

    def __set_content(self, content_type: str, shape) -> Image | TextBox:
        if content_type == config.IMAGE_KEY:
            return Image(
                coordinates=(shape.left, shape.top),
                size=(shape.width, shape.height),
                label=shape.text,
                path=""
            )
        if content_type == config.TEXTBOX_KEY:
            return TextBox(
                coordinates=(shape.left, shape.top),
                size=(shape.width, shape.height),
                label=shape.text,
                text=shape.text
            )

    def __get_contents(self, content_type: str, shape_type, slide_index=0):
        """特定のshape_typeを取得"""
        contents: list[Image | TextBox] = []
        shapes = self.prs.slides[slide_index].shapes

        for shape in shapes:
            if shape.shape_type != shape_type:
                continue
            if not shape.text:
                continue
            contents.append(self.__set_content(content_type, shape))

        return contents


class PresentationWriter():
    """
    プレゼンテーションを出力するクラス。

    prs: Presentation
        pptx.Presentationオブジェクト。
    slides: list[Slide]
        書き込みたい情報を持つSlideのリスト。
    slide_start_index: int
        スライドの最初のインデックス。
    layout
        スライドの作成に使うレイアウト情報。優先順位: 白紙 > タイトルとコンテンツ > 最後
    """
    def __init__(self, path: str, slides: list[Slide]) -> None:
        self.prs: Presentation = Presentation(path)
        self.slides: list[Slide] = slides
        # 0だと判定に使いにくいので1000で初期化する
        self.slide_start_index: int = 1000

        # get_by_nameは存在しない場合Noneを返す
        self.layout = self.prs.slide_layouts.get_by_name("白紙")
        if not self.layout:
            self.layout = self.prs.slide_layouts.get_by_name("タイトルとコンテンツ")
        if not self.layout:
            self.layout = self.prs.slide_layouts[-1]

    def set_empty_slides(self, base_slide_index=0):
        """空のスライドをセットする。
        置換しない図形は残しておきたいので、ベースとなるスライドに対して
        置換対象の図形をクリア → 複製 → 削除 という手順を取る。

        Parameters
        ----------
        base_slide_index : int, optional
            ベースとなるスライドのインデックス, by default 0
        """
        self.clear_pointing_images(base_slide_index)
        self.clear_textboxes(base_slide_index)

        for __ in self.slides:
            dup_index = self.duplicate_slide(base_slide_index)
            if self.slide_start_index == 1000:
                # テンプレートスライドは消すため-1する。
                self.slide_start_index = dup_index - 1

        self.remove_slide(base_slide_index)

    def add_images(self):
        """スライドに画像を貼り付ける。
        """
        current_index = self.slide_start_index

        for slide in self.slides:
            shapes = self.prs.slides[current_index].shapes
            for group in slide.contents:
                for image in group["image"]:
                    self.__add_image(image["path"], image["coordinates"], image["size"], shapes)
            current_index += 1

    def add_textboxes(self):
        """スライドにテキストボックスを貼り付ける。
        """
        current_index = self.slide_start_index

        for slide in self.slides:
            shapes = self.prs.slides[current_index].shapes
            for group in slide.contents:
                for textbox in group["textbox"]:
                    self.__add_textbox(textbox["text"], textbox["coordinates"], textbox["size"], shapes)
            current_index += 1

    def add_slide(self):
        """空のスライドを追加する。
        """
        return self.prs.slides.add_slide(self.layout)

    def duplicate_slide(self, base_slide_index=0) -> int:
        """スライドを複製する。
        参考 : https://stackoverflow.com/questions/50866634/python-pptx-copy-slide

        Parameters
        ----------
        base_slide_index : int, optional
            複製元スライドのインデックス, by default 0

        Returns
        -------
        int
            複製したスライドのインデックス。
        """
        base_slide = self.prs.slides[base_slide_index]
        duplicate_slide = self.add_slide()
        for shape in base_slide.shapes:
            el = shape.element
            new_el = deepcopy(el)
            duplicate_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
        return self.__get_slide_index(duplicate_slide)

    def remove_slide(self, slide_index=0) -> None:
        """スライドを削除する。

        Parameters
        ----------
        slide_index : int, optional
            対象スライドのインデックス, by default 0
        """
        xml_slides = self.prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[slide_index])

    def clear_pointing_images(self, slide_index=0) -> None:
        """スライドの画像を指す図形を削除する。

        Parameters
        ----------
        slide_index : int, optional
            対象スライドのインデックス, by default 0
        """
        self.__clear_shape(config.SHAPES["image"], slide_index)

    def clear_textboxes(self, slide_index=0) -> None:
        """スライドのテキストボックスをクリアする。

        Parameters
        ----------
        slide_index : int, optional
            対象スライドのインデックス, by default 0
        """
        self.__clear_shape(config.SHAPES["textbox"], slide_index)

    def __clear_shape(self, shape_type=config.SHAPES["image"], slide_index=0) -> None:
        """スライドからshapeを削除する。
        参考 : https://stackoverflow.com/questions/64700638/is-there-a-way-to-delete-a-shape-with-python-pptx

        Parameters
        ----------
        shape_type : pptx.enum.base.EnumValue, optional
            shapeオブジェクトのタイプ, by default config.SHAPES["image"]
        slide_index : int, optional
            対象スライドのインデックス, by default 0
        """
        shapes = self.prs.slides[slide_index].shapes
        for shape in shapes:
            if not shape.shape_type == shape_type:
                continue
            XML_reference = shape._sp
            XML_reference.getparent().remove(XML_reference)

    def __get_slide_index(self, slide) -> int:
        """スライドのインデックスを取得する。
        """
        return self.prs.slides.index(slide)

    def __add_image(
            self,
            img: str,
            coordinates: tuple[int, int],
            sizes: tuple[int, int],
            shapes
        ):
        """スライドに画像を貼り付ける。
        """
        shapes.add_picture(img, *coordinates, *sizes)

    def __add_textbox(
                self, text: str,
                coordinates: tuple[int, int],
                sizes: tuple[int, int],
                shapes
        ):
        textbox = shapes.add_textbox(*coordinates, *sizes)
        textbox.text_frame.text = text

    def save(self, dst_path: str) -> bool:
        """プレゼンテーションを保存する。

        Parameters
        ----------
        dst_path : str
            出力先のパス。

        Returns
        -------
        bool
            保存の成功/失敗を指す
        """
        try:
            self.prs.save(dst_path)
        except PermissionError:
            return False
        return True